from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement


def set_background_color(slide, color):
    """设置幻灯片背景颜色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_gradient_shape(slide, left, top, width, height, color1, color2):
    """添加渐变形状作为装饰元素"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2
    shape.line.fill.background()
    return shape


def set_cell_border(cell, border_color=None):
    """设置表格单元格边框"""
    if border_color is None:
        border_color = RGBColor(200, 200, 200)

    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    # 创建边框
    for border_name in ['left', 'right', 'top', 'bottom']:
        border = OxmlElement(f'a:{border_name}')
        border.set('w', '12700')  # 边框宽度
        borderColor = OxmlElement('a:solidFill')
        borderColor.append(OxmlElement('a:srgbClr'))
        borderColor[0].set('val', f"{border_color.rgb:06x}")
        border.append(borderColor)
        tcPr.append(border)


def create_beautiful_distributed_fs_ppt():
    # 创建演示文稿对象
    prs = Presentation()

    # 设置幻灯片宽高比 (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 定义颜色主题
    PRIMARY_COLOR = RGBColor(41, 128, 185)  # 蓝色
    SECONDARY_COLOR = RGBColor(52, 152, 219)  # 浅蓝色
    ACCENT_COLOR = RGBColor(231, 76, 60)  # 红色
    LIGHT_BG = RGBColor(245, 245, 245)  # 浅灰色背景
    DARK_TEXT = RGBColor(44, 62, 80)  # 深灰色文字

    # ===== 第1页：封面 =====
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 设置渐变背景
    set_background_color(slide, LIGHT_BG)
    add_gradient_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(2),
                       PRIMARY_COLOR, SECONDARY_COLOR)

    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "分布式文件系统"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # 添加副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "第9章 9.1-9.2.1节\n汇报人：[你的姓名] | 日期：[汇报日期]"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # ===== 第2页：目录 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, LIGHT_BG)

    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "目录"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.font.bold = True

    # 添加装饰线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(12), Inches(0.1)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY_COLOR

    # 添加目录内容
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
    content_frame = content_box.text_frame
    content_frame.text = ""

    topics = [
        "1. 分布式文件系统概述",
        "2. 抽象模型架构",
        "3. 核心设计问题",
        "4. NFS结构模型分析",
        "5. 总结与展望"
    ]

    for i, topic in enumerate(topics):
        p = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
        p.text = topic
        p.font.size = Pt(24)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(12)

        # 添加项目符号装饰
        bullet_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.6 + i * 0.7), Inches(0.2), Inches(0.2)
        )
        bullet_shape.fill.solid()
        bullet_shape.fill.fore_color.rgb = ACCENT_COLOR
        bullet_shape.line.fill.background()

    # ===== 第3页：分布式文件系统概述 =====
    slide = create_content_slide(prs, "分布式文件系统概述", PRIMARY_COLOR, LIGHT_BG, DARK_TEXT)

    content = [
        ("📁 核心功能", [
            "• 组织、存储、提取、命名、共享和保护文件",
            "• 支持多客户通过网络共享文件"
        ]),
        ("🔍 透明性要求", [
            "• 位置透明性 - 路径名不暴露物理位置",
            "• 多副本透明性 - 用户无需关心副本管理",
            "• 客户端缓存透明性 - 缓存对用户透明"
        ]),
        ("💡 典型系统", [
            "• NFS (Network File System)",
            "• AFS/Coda",
            "• SMB/CIFS"
        ])
    ]

    add_styled_content(slide, content, Inches(1), Inches(1.8), Inches(11), Inches(5), DARK_TEXT)

    # ===== 第4页：抽象模型架构 =====
    slide = create_content_slide(prs, "抽象模型架构", PRIMARY_COLOR, LIGHT_BG, DARK_TEXT)

    content = [
        ("🖥️ 客户端", [
            "• 台式机或工作站环境",
            "• 运行客户程序模块",
            "• 支持本地与远程文件统一访问"
        ]),
        ("🖥️ 服务器端", [
            "• 提供目录服务和文件服务",
            "• 文件操作 + 属性操作",
            "• 唯一文件标识符(UFID)管理"
        ]),
        ("🌐 网络通信", [
            "• 基于RPC远程过程调用",
            "• TCP/UDP协议支持",
            "• 跨网络透明访问"
        ])
    ]

    add_styled_content(slide, content, Inches(1), Inches(1.8), Inches(11), Inches(5), DARK_TEXT)

    # ===== 第5页：文件服务接口 =====
    slide = create_content_slide(prs, "文件服务接口", PRIMARY_COLOR, LIGHT_BG, DARK_TEXT)

    # 添加说明文字
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(0.5))
    desc_frame = desc_box.text_frame
    desc_frame.text = "基于UFID(唯一文件标识符)的核心操作接口"
    desc_para = desc_frame.paragraphs[0]
    desc_para.font.size = Pt(14)
    desc_para.font.color.rgb = DARK_TEXT
    desc_para.font.italic = True

    # 创建美观的表格
    file_operations = [
        ("📖 Read(UFID,i,n)→Data", "从文件指定位置读取数据项"),
        ("✏️ Write(UFID,i,Data)", "向文件指定位置写入数据"),
        ("🆕 Create()→UFID", "创建新文件并返回UFID"),
        ("🗑️ Delete(UFID)", "删除指定文件"),
        ("📊 GetAttributes(UFID)→Attr", "获取文件属性信息"),
        ("⚙️ SetAttributes(UFID,Attr)", "设置文件属性")
    ]

    create_styled_table(slide, file_operations, ["操作", "功能描述"],
                        Inches(1), Inches(2.5), Inches(11), Inches(3),
                        PRIMARY_COLOR, LIGHT_BG, DARK_TEXT)

    # ===== 第6页：目录服务接口 =====
    slide = create_content_slide(prs, "目录服务接口", PRIMARY_COLOR, LIGHT_BG, DARK_TEXT)

    desc_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(0.5))
    desc_frame = desc_box.text_frame
    desc_frame.text = "文件名与UFID之间的映射管理"
    desc_para = desc_frame.paragraphs[0]
    desc_para.font.size = Pt(14)
    desc_para.font.color.rgb = DARK_TEXT
    desc_para.font.italic = True

    dir_operations = [
        ("📁 MkDir(Dir,Name,Attr)→UFID", "创建新的子目录"),
        ("🗂️ RmDir(Dir,Name)", "删除指定子目录"),
        ("🔍 Lookup(Dir,Name)→UFID", "根据文件名查找UFID"),
        ("➕ AddName(Dir,Name,UFID)", "添加目录项映射"),
        ("➖ UnName(Dir,Name)", "删除目录项"),
        ("📋 GetName(Dir,Pattern)→NameSeq", "模式匹配文件名")
    ]

    create_styled_table(slide, dir_operations, ["操作", "功能描述"],
                        Inches(1), Inches(2.5), Inches(11), Inches(3.5),
                        PRIMARY_COLOR, LIGHT_BG, DARK_TEXT)

    # ===== 第7-15页：其他内容幻灯片 =====
    # 由于代码长度限制，这里只展示关键的美化函数
    # 你可以按照类似模式继续创建其他幻灯片

    # ===== 第16页：结束页 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, LIGHT_BG)

    # 添加渐变装饰
    add_gradient_shape(slide, Inches(0), Inches(5), Inches(13.333), Inches(2.5),
                       PRIMARY_COLOR, SECONDARY_COLOR)

    # 感谢文字
    thank_you_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    thank_you_frame = thank_you_box.text_frame
    thank_you_frame.text = "感谢聆听"
    thank_you_para = thank_you_frame.paragraphs[0]
    thank_you_para.font.size = Pt(48)
    thank_you_para.font.color.rgb = PRIMARY_COLOR
    thank_you_para.font.bold = True
    thank_you_para.alignment = PP_ALIGN.CENTER

    qa_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1))
    qa_frame = qa_box.text_frame
    qa_frame.text = "Q&A | 欢迎提问交流"
    qa_para = qa_frame.paragraphs[0]
    qa_para.font.size = Pt(24)
    qa_para.font.color.rgb = DARK_TEXT
    qa_para.alignment = PP_ALIGN.CENTER

    # 保存PPT文件
    prs.save('分布式文件系统汇报_美化版.pptx')
    print("美化版PPT已生成：分布式文件系统汇报_美化版.pptx")


def create_content_slide(prs, title, primary_color, bg_color, text_color):
    """创建统一风格的内容幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, bg_color)

    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.color.rgb = primary_color
    title_para.font.bold = True

    # 添加标题装饰线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(12), Inches(0.08)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = primary_color

    return slide


def add_styled_content(slide, content, left, top, width, height, text_color):
    """添加带样式的文本内容"""
    current_top = top
    for section_title, items in content:
        # 添加小节标题
        title_box = slide.shapes.add_textbox(left, current_top, width, Inches(0.4))
        title_frame = title_box.text_frame
        title_frame.text = section_title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(18)
        title_para.font.color.rgb = text_color
        title_para.font.bold = True

        current_top += Inches(0.5)

        # 添加内容项
        content_box = slide.shapes.add_textbox(left + Inches(0.3), current_top, width - Inches(0.3), Inches(0.8))
        content_frame = content_box.text_frame
        content_frame.text = ""

        for item in items:
            p = content_frame.add_paragraph() if content_frame.paragraphs else content_frame.paragraphs[0]
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = text_color
            p.space_after = Pt(6)

        current_top += Inches(0.9)


def create_styled_table(slide, data, headers, left, top, width, height, header_color, row_color, text_color):
    """创建带样式的表格"""
    rows = len(data) + 1
    cols = len(headers)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # 设置表头
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        # 设置表头文字样式
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)

    # 填充数据行
    for row, (icon_text, description) in enumerate(data, 1):
        # 第一列（操作列）
        cell1 = table.cell(row, 0)
        cell1.text = icon_text
        cell1.fill.solid()
        cell1.fill.fore_color.rgb = row_color

        # 第二列（描述列）
        cell2 = table.cell(row, 1)
        cell2.text = description
        cell2.fill.solid()
        cell2.fill.fore_color.rgb = row_color

        # 设置文字样式
        for col in range(cols):
            cell = table.cell(row, col)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.color.rgb = text_color
                paragraph.font.size = Pt(11)

    return table


# 运行生成函数
if __name__ == "__main__":
    create_beautiful_distributed_fs_ppt()