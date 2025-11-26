from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def create_distributed_file_system_ppt():
    # 创建演示文稿对象
    prs = Presentation()

    # 设置幻灯片宽高比 (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===== 第1页：封面 =====
    slide_layout = prs.slide_layouts[0]  # 标题幻灯片布局
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "分布式文件系统"
    subtitle.text = "第9章 9.1-9.2.1节\n\n汇报人：[你的姓名]\n日期：[汇报日期]"

    # ===== 第2页：目录 =====
    slide_layout = prs.slide_layouts[1]  # 标题和内容布局
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "目录"
    tf = content.text_frame
    tf.text = "主要内容"

    p = tf.add_paragraph()
    p.text = "1. 分布式文件系统概述"

    p = tf.add_paragraph()
    p.text = "2. 抽象模型"

    p = tf.add_paragraph()
    p.text = "3. 设计问题"

    p = tf.add_paragraph()
    p.text = "4. NFS结构模型"

    p = tf.add_paragraph()
    p.text = "5. 总结"

    # ===== 第3页：分布式文件系统概述 =====
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "分布式文件系统概述"
    tf = content.text_frame
    tf.text = "核心功能："

    p = tf.add_paragraph()
    p.text = "• 组织、存储、提取、命名、共享和保护文件"

    p = tf.add_paragraph()
    p.text = "• 支持多客户通过网络共享文件"

    p = tf.add_paragraph()
    p.text = "\n透明性要求："

    p = tf.add_paragraph()
    p.text = "• 位置透明性"

    p = tf.add_paragraph()
    p.text = "• 多副本透明性"

    p = tf.add_paragraph()
    p.text = "• 客户端缓存透明性"

    p = tf.add_paragraph()
    p.text = "\n典型系统：NFS, AFS/Coda, SMB/CIFS"

    # ===== 第4页：抽象模型架构 =====
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "抽象模型架构"
    tf = content.text_frame
    tf.text = "客户端："

    p = tf.add_paragraph()
    p.text = "• 台式机或工作站"

    p = tf.add_paragraph()
    p.text = "• 运行客户程序"

    p = tf.add_paragraph()
    p.text = "• 安装客户模块访问本地和远程文件"

    p = tf.add_paragraph()
    p.text = "\n服务器端："

    p = tf.add_paragraph()
    p.text = "• 提供目录服务和文件服务"

    p = tf.add_paragraph()
    p.text = "• 文件服务包括文件操作和属性操作"

    p = tf.add_paragraph()
    p.text = "\n网络：通过RPC实现远程操作"

    # ===== 第5页：文件服务接口 =====
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "文件服务接口"
    tf = content.text_frame
    tf.text = "文件服务操作："

    # 添加表格
    rows, cols = 7, 2
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # 设置表头
    table.cell(0, 0).text = "操作"
    table.cell(0, 1).text = "功能描述"

    # 填充表格内容
    file_operations = [
        ("Read(UFID,i,n)→Data", "从文件指定位置读取数据"),
        ("Write(UFID,i,Data)", "向文件指定位置写入数据"),
        ("Create()→UFID", "创建新文件并返回UFID"),
        ("Delete(UFID)", "删除指定文件"),
        ("GetAttributes(UFID)→Attr", "获取文件属性"),
        ("SetAttributes(UFID,Attr)", "设置文件属性")
    ]

    for i, (op, desc) in enumerate(file_operations, 1):
        table.cell(i, 0).text = op
        table.cell(i, 1).text = desc

    # ===== 第6页：目录服务接口 =====
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "目录服务接口"
    tf = content.text_frame
    tf.text = "目录服务操作："

    # 添加表格
    rows, cols = 7, 2
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # 设置表头
    table.cell(0, 0).text = "操作"
    table.cell(0, 1).text = "功能描述"

    # 填充表格内容
    dir_operations = [
        ("MkDir(Dir,Name,Attr)→UFID", "创建子目录"),
        ("RmDir(Dir,Name)", "删除子目录"),
        ("Lookup(Dir,Name)→UFID", "查找文件UFID"),
        ("AddName(Dir,Name,UFID)", "添加目录项"),
        ("UnName(Dir,Name)", "删除目录项"),
        ("GetName(Dir,Pattern)→NameSeq", "匹配文件名")
    ]

    for i, (op, desc) in enumerate(dir_operations, 1):
        table.cell(i, 0).text = op
        table.cell(i, 1).text = desc

    # ===== 第7页：设计问题 - 文件使用模式 =====
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "设计问题 - 文件使用模式"
    tf = content.text_frame
    tf.text = "调查发现："

    p = tf.add_paragraph()
    p.text = "• 大多数文件长度小于10KB"

    p = tf.add_paragraph()
    p.text = "• 大多数文件生命周期短"

    p = tf.add_paragraph()
    p.text = "• 普通数据文件可能被共享"

    p = tf.add_paragraph()
    p.text = "\n设计策略："

    p = tf.add_paragraph()
    p.text = "• 支持大文件传输"

    p = tf.add_paragraph()
    p.text = "• 本地保存临时文件"

    p = tf.add_paragraph()
    p.text = "• 选择合适的共享语义"

    # ===== 第8页：设计问题 - 命名与名字解析 =====
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "设计问题 - 命名与名字解析"
    tf = content.text_frame
    tf.text = "关键概念："

    p = tf.add_paragraph()
    p.text = "• 位置透明性：路径名不暴露物理位置"

    p = tf.add_paragraph()
    p.text = "• 位置无关性：文件移动不影响路径名"

    p = tf.add_paragraph()
    p.text = "\n名字空间构造方式："

    p = tf.add_paragraph()
    p.text = "1. 机器名+路径名：/machine/path"

    p = tf.add_paragraph()
    p.text = "2. 挂载远程文件到本地目录"

    p = tf.add_paragraph()
    p.text = "3. 单一全局名字空间"

    # ===== 第9页：设计问题 - 访问模型 =====
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "设计问题 - 访问模型"
    tf = content.text_frame
    tf.text = "远程访问模型："

    p = tf.add_paragraph()
    p.text = "• 服务器实现文件操作"

    p = tf.add_paragraph()
    p.text = "• 客户端无需大存储空间"

    p = tf.add_paragraph()
    p.text = "\n上载/下载模型："

    p = tf.add_paragraph()
    p.text = "• 客户端下载整个文件到本地"

    p = tf.add_paragraph()
    p.text = "• 操作完成后上传回服务器"

    p = tf.add_paragraph()
    p.text = "• 概念简单但可能浪费带宽"

    # ===== 第10页：设计问题 - 缓存策略 =====
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "设计问题 - 缓存策略"
    tf = content.text_frame
    tf.text = "缓存位置："

    p = tf.add_paragraph()
    p.text = "• 服务器主存"

    p = tf.add_paragraph()
    p.text = "• 客户端磁盘"

    p = tf.add_paragraph()
    p.text = "• 客户端主存"

    p = tf.add_paragraph()
    p.text = "\n一致性解决方案："

    p = tf.add_paragraph()
    p.text = "• 直写(Write-Through)"

    p = tf.add_paragraph()
    p.text = "• 推后写(Delayed Write)"

    p = tf.add_paragraph()
    p.text = "• 关闭写(Write-on-Close)"

    p = tf.add_paragraph()
    p.text = "• 集中控制(Centralized Control)"

    # ===== 第11页：设计问题 - 文件共享语义 =====
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "设计问题 - 文件共享语义"
    tf = content.text_frame
    tf.text = "UNIX语义："

    p = tf.add_paragraph()
    p.text = "• 顺序一致性"

    p = tf.add_paragraph()
    p.text = "• 写后读得到最新结果"

    p = tf.add_paragraph()
    p.text = "\n会话语义："

    p = tf.add_paragraph()
    p.text = "• 修改仅在文件关闭后对其他客户端可见"

    p = tf.add_paragraph()
    p.text = "\n不修改共享文件语义："

    p = tf.add_paragraph()
    p.text = "• 文件只读，修改需创建新文件"

    p = tf.add_paragraph()
    p.text = "\n事务语义："

    p = tf.add_paragraph()
    p.text = "• 原子性操作保证"

    # ===== 第12页：设计问题 - 容错与状态 =====
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "设计问题 - 容错与状态"
    tf = content.text_frame
    tf.text = "无状态服务："

    p = tf.add_paragraph()
    p.text = "• 请求完全自包含"

    p = tf.add_paragraph()
    p.text = "• 容错性好"

    p = tf.add_paragraph()
    p.text = "• 不维护客户端状态"

    p = tf.add_paragraph()
    p.text = "\n有状态服务："

    p = tf.add_paragraph()
    p.text = "• 维护客户端连接信息"

    p = tf.add_paragraph()
    p.text = "• 支持文件锁定"

    p = tf.add_paragraph()
    p.text = "• 性能较好"

    p = tf.add_paragraph()
    p.text = "\n多副本：提高可用性和性能"
    p = tf.add_paragraph()
    p.text = "断开操作：支持移动客户端离线工作"

    # ===== 第13页：NFS结构模型 =====
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "NFS结构模型"
    tf = content.text_frame
    tf.text = "核心组件："

    p = tf.add_paragraph()
    p.text = "• 虚拟文件系统(VFS)：隐藏不同文件系统差异"

    p = tf.add_paragraph()
    p.text = "• vfs对象：代表一个文件系统"

    p = tf.add_paragraph()
    p.text = "• vnode对象：代表文件或目录"

    p = tf.add_paragraph()
    p.text = "\n通信机制："

    p = tf.add_paragraph()
    p.text = "• 基于RPC的远程文件访问"

    p = tf.add_paragraph()
    p.text = "• 支持硬连接和符号连接"

    # ===== 第14页：NFS版本对比 =====
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "NFS版本对比"
    tf = content.text_frame
    tf.text = "NFS操作对比："

    # 添加表格
    rows, cols = 5, 4
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(10)
    height = Inches(3)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # 设置表头
    headers = ["操作", "NFSv3", "NFSv4", "说明"]
    for i, header in enumerate(headers):
        table.cell(0, i).text = header

    # 填充表格内容
    nfs_comparison = [
        ("create", "✓", "✓", "创建文件"),
        ("open/close", "✗", "✓", "NFSv4新增"),
        ("lookup", "不跨挂载点", "可跨挂载点", "解析方式不同")
    ]

    for i, (op, v3, v4, desc) in enumerate(nfs_comparison, 1):
        table.cell(i, 0).text = op
        table.cell(i, 1).text = v3
        table.cell(i, 2).text = v4
        table.cell(i, 3).text = desc

    # 添加NFSv4改进说明
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(10), Inches(1.5))
    tf = textbox.text_frame
    tf.text = "NFSv4改进："
    p = tf.add_paragraph()
    p.text = "• 增强安全性"
    p = tf.add_paragraph()
    p.text = "• 支持有状态服务"
    p = tf.add_paragraph()
    p.text = "• 改进名字解析"

    # ===== 第15页：总结 =====
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "总结"
    tf = content.text_frame
    tf.text = "分布式文件系统核心要求："

    p = tf.add_paragraph()
    p.text = "• 透明性、一致性、容错性、性能"

    p = tf.add_paragraph()
    p.text = "\n关键设计问题："

    p = tf.add_paragraph()
    p.text = "• 命名、缓存、共享语义、状态管理"

    p = tf.add_paragraph()
    p.text = "\nNFS演进："

    p = tf.add_paragraph()
    p.text = "• 从无状态到有状态"

    p = tf.add_paragraph()
    p.text = "• 功能不断增强"

    p = tf.add_paragraph()
    p.text = "• 成为分布式文件系统事实标准"

    # ===== 第16页：Q&A =====
    slide_layout = prs.slide_layouts[0]  # 标题幻灯片布局
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Q&A"
    subtitle.text = "谢谢聆听！\n欢迎提问！"

    # 保存PPT文件
    prs.save('分布式文件系统汇报.pptx')
    print("PPT已生成：分布式文件系统汇报.pptx")


# 运行生成函数
if __name__ == "__main__":
    create_distributed_file_system_ppt()